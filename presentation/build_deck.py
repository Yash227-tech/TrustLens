"""Generate TrustLens_Final.pptx — dark neon-green theme matching the web app
(mockup #3: near-black green bg + neon-mint accents, glass cards, NO gradients).

Includes a detailed end-to-end workflow diagram + step-by-step process slides."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- theme tokens (from frontend/src/index.css .dark) ----
BG      = RGBColor(0x0B, 0x14, 0x10)
CARD    = RGBColor(0x12, 0x20, 0x1A)
CARD2   = RGBColor(0x16, 0x27, 0x20)
BORDER  = RGBColor(0x25, 0x38, 0x2F)
NEON    = RGBColor(0x10, 0xE8, 0x9C)
TEXT    = RGBColor(0xEA, 0xF6, 0xF0)
MUTED   = RGBColor(0x9D, 0xB0, 0xA8)
AMBER   = RGBColor(0xF4, 0xA5, 0x2A)
RED     = RGBColor(0xEF, 0x44, 0x44)
FONT    = "Segoe UI"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]
_N = [0]


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background(); r.shadow.inherit = False
    return s


def rect(s, x, y, w, h, fill=None, line=None, radius=0.06, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sh = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(1)
    sh.shadow.inherit = False
    try:
        if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
            sh.adjustments[0] = radius
    except Exception:
        pass
    return sh


def text(s, x, y, w, h, runs, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, spacing=1.0, font=FONT):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [[(runs, color, bold, size)]]
    elif runs and not isinstance(runs[0], list):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing; p.space_after = Pt(3)
        if isinstance(para, str):
            para = [(para, color, bold, size)]
        for tup in para:
            txt, c, b, sz = (tup + (size,))[:4] if len(tup) < 4 else tup
            r = p.add_run(); r.text = txt
            r.font.size = Pt(sz); r.font.bold = b; r.font.color.rgb = c; r.font.name = font
    return tb


def header(s, kicker, title):
    _N[0] += 1
    rect(s, 0.7, 0.42, 0.17, 0.17, fill=NEON, radius=0.25)
    text(s, 0.95, 0.4, 4, 0.3, [("Trust", TEXT, True, 14), ("Lens", NEON, True, 14)])
    text(s, 0.7, 1.02, 12, 0.3, kicker.upper(), size=12.5, color=NEON, bold=True)
    text(s, 0.7, 1.34, 12, 0.9, title, size=27, color=TEXT, bold=True)
    text(s, 11.2, 7.02, 1.4, 0.3, f"{_N[0]:02d}", size=11, color=MUTED, align=PP_ALIGN.RIGHT)
    text(s, 0.7, 7.02, 8, 0.3, "TrustLens — AI Document-Fraud Detection", size=10, color=MUTED)


def bullets(s, x, y, w, items, size=15.5, gap=0.62, color=TEXT):
    for i, it in enumerate(items):
        head, sub = (it if isinstance(it, tuple) else (it, None))
        text(s, x, y + i * gap, w, gap, [(("▸  "), NEON, True, size), (head, color, False, size)])
        if sub:
            text(s, x + 0.32, y + i * gap + 0.27, w - 0.32, gap, sub, size=12, color=MUTED)


def card(s, x, y, w, h, fill=CARD):
    return rect(s, x, y, w, h, fill=fill, line=BORDER, radius=0.07)


def metric(s, x, y, w, h, value, label, vcolor=NEON):
    card(s, x, y, w, h)
    text(s, x + 0.25, y + 0.22, w - 0.5, 0.7, value, size=30, color=vcolor, bold=True)
    text(s, x + 0.27, y + h - 0.78, w - 0.5, 0.7, label, size=12.5, color=MUTED)


def badge(s, x, y, n, d=0.42):
    rect(s, x, y, d, d, fill=NEON, radius=0.25)
    text(s, x, y + 0.02, d, d - 0.02, str(n), size=16, color=BG, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def arrow(s, x, y, w=0.45, color=NEON):
    text(s, x, y, w, 0.5, "→", size=20, color=color, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ===================== 1 · TITLE =====================
s = slide()
rect(s, 0.9, 1.7, 0.55, 0.55, fill=NEON, radius=0.28)
text(s, 1.65, 1.62, 10, 1.2, [("Trust", TEXT, True, 60), ("Lens", NEON, True, 60)])
text(s, 0.92, 2.95, 11.5, 0.7, "AI Document-Fraud Detection for Bank Underwriting", size=24, color=TEXT, bold=True)
text(s, 0.92, 3.62, 11.5, 0.5,
     [("On-premise", NEON, True, 15), ("  ·  Explainable  ·  Human-in-the-loop  ·  RBI / DPDP-friendly", MUTED, False, 15)])
for i, (v, l, c) in enumerate([
        ("98.6%", "Fraud caught (benchmark)", NEON), ("26", "Document types", TEXT),
        ("0", "Customer-data egress", NEON), ("<5s", "Per-document analysis", TEXT)]):
    metric(s, 0.9 + i * 3.1, 4.7, 2.85, 1.5, v, l, vcolor=c)
text(s, 0.7, 6.95, 12, 0.3, "Canara Bank Hackathon · Final model", size=11, color=MUTED)

# ===================== 2 · PROBLEM =====================
s = slide(); header(s, "The problem", "Forged loan documents slip through manual review")
bullets(s, 0.75, 2.5, 7.3, [
    ("Loan files are bundles of forgeable documents", "Edited PDFs, pasted signatures, fake stamps, fabricated income & identity."),
    ("Manual verification is slow and inconsistent", "3–7 days per file; reviewers fatigue and miss sophisticated edits."),
    ("Sophisticated forgeries beat the human eye", "A digit changed in Word, a swapped ID photo, a real PAN with the wrong name."),
    ("Banks need speed AND a defensible audit trail", "Every decision must be explainable and logged for RBI compliance."),
], gap=1.0)
card(s, 8.5, 2.5, 4.1, 3.7)
text(s, 8.8, 2.75, 3.6, 0.5, "Cost of a miss", size=13, color=NEON, bold=True)
text(s, 8.8, 3.3, 3.6, 2.6, [
    [("Financial loss", TEXT, True, 15)], [("Fraudulent disbursement on a fabricated file.", MUTED, False, 12.5)],
    [("Regulatory risk", TEXT, True, 15)], [("No immutable record of why a file was approved.", MUTED, False, 12.5)],
    [("Slow genuine customers", TEXT, True, 15)], [("Honest applicants wait days behind manual checks.", MUTED, False, 12.5)],
], spacing=1.05)

# ===================== 3 · SOLUTION =====================
s = slide(); header(s, "The solution", "Trust Score, visual evidence & risk routing — in seconds")
cols = [("Analyse", "10 forensic + AI checks on every page — metadata, fonts, ELA, copy-move, stamps, faces."),
        ("Score & explain", "An XGBoost Trust Score (0–100) with SHAP attribution and a plain-language LLM report."),
        ("Route by risk", "GREEN fast-track · YELLOW review · RED escalate. It never auto-rejects."),
        ("Log immutably", "Every decision written to an append-only audit trail for compliance.")]
for i, (t, d) in enumerate(cols):
    x = 0.75 + i * 3.07
    card(s, x, 2.6, 2.95, 3.2)
    rect(s, x + 0.25, 2.85, 0.45, 0.45, fill=NEON, radius=0.25)
    text(s, x + 0.32, 2.9, 0.4, 0.4, str(i + 1), size=18, color=BG, bold=True)
    text(s, x + 0.25, 3.5, 2.45, 0.5, t, size=17, color=TEXT, bold=True)
    text(s, x + 0.25, 4.05, 2.45, 1.6, d, size=12.5, color=MUTED, spacing=1.1)
text(s, 0.75, 6.1, 12, 0.5,
     [("Bank-safe by design: ", NEON, True, 14), ("favours review over rejection — the human is faster and more confident, never replaced.", TEXT, False, 14)])

# ===================== 4 · PIPELINE (at a glance) =====================
s = slide(); header(s, "How it works — at a glance", "One upload → a full forensic + AI pipeline")
stages = ["Upload\n(PDF / image\n/ DOCX)", "Extract & classify\n(OCR · LayoutLMv3\n· NER + regex)",
          "Forensics\n(7 signals +\ncritical checks)", "Cross-source\n(DigiLocker · GSTN\n· Udyam · face)",
          "Trust Score\n(XGBoost + SHAP\n+ LLM report)", "Tier & audit\n(GREEN/YELLOW/RED\n· append-only log)"]
bw, bh, by = 1.92, 1.7, 3.0
gap = (12.0 - bw) / (len(stages) - 1)
for i, st in enumerate(stages):
    x = 0.7 + i * gap
    card(s, x, by, bw, bh, fill=CARD2 if i % 2 else CARD)
    tl, *rest = st.split("\n")
    text(s, x + 0.12, by + 0.2, bw - 0.24, bh - 0.3,
         [[(tl, NEON, True, 13)]] + [[(r, MUTED, False, 10.5)] for r in rest],
         spacing=1.0, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < len(stages) - 1:
        arrow(s, x + bw - 0.02, by + bh / 2 - 0.25, gap - bw + 0.04)
text(s, 0.75, 5.4, 12, 0.6,
     [("Runs entirely on-premise (Docker + RTX 4060). ", TEXT, True, 14),
      ("No customer data ever leaves the bank.", NEON, True, 14)])

# ===================== 5 · WORKFLOW DIAGRAM (top-to-bottom flowchart) =====================
s = slide(); header(s, "End-to-end workflow", "Four stages, a branch, two feedback loops")


def glyph(s, x, y, w, h, ch, size=18, color=NEON):
    text(s, x, y, w, h, ch, size=size, color=color, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


SPW = 6.7
sx = (SW - SPW) / 2          # centred vertical spine
spine = [
    ("1", "INGEST", "Upload PDF / image / DOCX  →  FastAPI · Celery + Redis · stored to MinIO"),
    ("2", "UNDERSTAND", "Normalise + OCR (en·hi·gu) · Classify: LayoutLMv3 + keyword (26 types) · Entities: regex + Verhoeff + NER"),
    ("3", "ANALYSE", "7 forensic signals · ELA · ManTraNet · stamp · font · photo-region  ·  Cross-source verify"),
    ("4", "SCORE", "XGBoost Trust Score 0–100 + SHAP attribution  ·  local LLM evidence report"),
]
y0, bh, stp = 2.3, 0.72, 0.92
for i, (n, t, d) in enumerate(spine):
    y = y0 + i * stp
    card(s, sx, y, SPW, bh)
    badge(s, sx + 0.16, y + 0.16, n, d=0.4)
    text(s, sx + 0.68, y, 1.95, bh, t, size=13.5, color=NEON, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    text(s, sx + 2.55, y + 0.05, SPW - 2.7, bh - 0.1, d, size=10.8, color=MUTED, anchor=MSO_ANCHOR.MIDDLE, spacing=1.0)
    if i < len(spine) - 1:
        glyph(s, sx + SPW / 2 - 0.25, y + bh - 0.02, 0.5, stp - bh + 0.04, "↓")
# branch into the three risk outcomes
ylast = y0 + (len(spine) - 1) * stp + bh
glyph(s, sx + SPW / 2 - 0.25, ylast - 0.02, 0.5, 0.28, "↓")
chips = [("GREEN", "Fast-track", NEON), ("YELLOW", "Review", AMBER), ("RED", "Escalate", RED)]
cw, cy = 2.1, ylast + 0.28
for i, (t, sub, c) in enumerate(chips):
    x = sx + i * (cw + 0.2)
    card(s, x, cy, cw, 0.64)
    rect(s, x, cy, cw, 0.1, fill=c, radius=0.0, shape=MSO_SHAPE.RECTANGLE)
    text(s, x + 0.2, cy, cw - 0.35, 0.64, [(t + "   ", c, True, 13.5), (sub, MUTED, False, 11)], anchor=MSO_ANCHOR.MIDDLE)
text(s, sx, cy + 0.72, SPW, 0.3, [("All outcomes  →  ", MUTED, False, 11), ("append-only audit log", NEON, True, 11)], align=PP_ALIGN.CENTER)

# feedback loops flanking the spine
ly, lh = 3.45, 1.45
# Case loop (left) → into the spine
card(s, 0.7, ly, 2.4, lh, fill=CARD2)
text(s, 0.88, ly + 0.16, 2.04, 0.4, "CASE LOOP", size=12, color=NEON, bold=True)
text(s, 0.88, ly + 0.58, 2.04, lh - 0.66, "Identity + face-match across the whole document bundle feed the tier.", size=10.8, color=MUTED, spacing=1.1)
glyph(s, 3.1, ly + lh / 2 - 0.2, sx - 3.1, 0.4, "→", size=16)
# Learning loop (right) ← from the decision back to the models
rx = 10.23
card(s, rx, ly, 12.63 - rx, lh, fill=CARD2)
text(s, rx + 0.18, ly + 0.16, 2.04, 0.4, "LEARNING LOOP", size=12, color=NEON, bold=True)
text(s, rx + 0.18, ly + 0.58, 2.04, lh - 0.66, "Underwriter's genuine/fraud verdict → flywheel → the models retrain.", size=10.8, color=MUTED, spacing=1.1)
glyph(s, sx + SPW, ly + lh / 2 - 0.2, rx - (sx + SPW), 0.4, "←", size=16)

# ===================== 6 · WORKFLOW DETAIL part 1 =====================
def step_grid(s, steps):
    cw, ch = 5.95, 1.28
    for i, (num, t, d) in enumerate(steps):
        col, r = divmod(i, 3)
        x = 0.7 + col * (cw + 0.2); y = 2.4 + r * (ch + 0.13)
        card(s, x, y, cw, ch)
        badge(s, x + 0.22, y + 0.22, num, d=0.4)
        text(s, x + 0.74, y + 0.2, cw - 0.95, 0.4, t, size=14, color=NEON, bold=True)
        text(s, x + 0.74, y + 0.6, cw - 0.95, 0.6, d, size=11.5, color=MUTED, spacing=1.08)


s = slide(); header(s, "Workflow in detail · 1 of 2", "From upload to forensic analysis")
step_grid(s, [
    ("1", "Ingestion", "Upload (PDF/JPG/PNG/DOCX) → FastAPI /api/analyze; file saved to disk + MinIO, a job_id is issued and a Celery task is queued via Redis (async, non-blocking)."),
    ("2", "Normalisation", "DOCX→PDF via LibreOffice; PDF pages rendered with PyMuPDF; images decoded. Detects born-digital vs scanned — which drives bank-safe calibration."),
    ("3", "Text extraction", "Born-digital text layer when present, else Tesseract OCR in English + Hindi + Gujarati. The method (scan vs digital) is recorded for downstream logic."),
    ("4", "Classification (dual)", "Keyword DocSpec anchors (authoritative, explainable) + fine-tuned LayoutLMv3 (visual second opinion) → one of 26 doc types; agreement/disagreement flagged."),
    ("5", "Entity extraction", "Regex for PAN, Aadhaar (+UIDAI Verhoeff), GSTIN, IFSC, Udyam, ITR-ack, passport-MRZ + spaCy NER (names/orgs). YOLO field detectors recover IDs from poor photos."),
    ("6", "Forensic engine", "Seven scored signals (PDF metadata, font/spacing, signature-ELA, stamp YOLO+SIFT, bank-balance, ELA, ManTraNet) + localised criticals (photo-region tamper, Udyam QR)."),
])

# ===================== 7 · WORKFLOW DETAIL part 2 =====================
s = slide(); header(s, "Workflow in detail · 2 of 2", "From verification to decision, audit & learning")
step_grid(s, [
    ("7", "Cross-source verify", "Parallel calls to DigiLocker, GSTN, Account Aggregator, Income-Tax e-Filing & Udyam; the authoritative name must match. Fabricated Aadhaar fails the Verhoeff check."),
    ("8", "Trust scoring", "The 7-feature vector → XGBoost → Trust Score 0–100, with SHAP per-feature attribution showing exactly which signal drove the result. Weighted fallback if model absent."),
    ("9", "Tier decision", "GREEN (≥85, no critical) · YELLOW (review) · RED (<50 or a critical forgery indicator). Review-only flags never auto-reject a genuine doc. The human always decides."),
    ("10", "Evidence report", "A local Llama 3.1 8B model writes a plain-language explanation; ELA/ManTraNet heatmaps are saved to MinIO and shown in the evidence viewer."),
    ("11", "Persist + audit", "Result is stored and an append-only AuditLog row freezes the full report (signals, SHAP, entities, heatmap) — immutable, for RBI compliance."),
    ("12", "Case + flywheel", "Across a bundle: identity + face-match consistency (conflict → RED/review). The underwriter's confirmed verdict feeds active-learning retraining."),
])

# ===================== 8 · MODEL STACK =====================
s = slide(); header(s, "Multi-model AI", "An ensemble — each model does what it is best at")
rows = [("LayoutLMv3 (fine-tuned)", "Document classification — 26 types", "99.3% val · 97.9% real"),
        ("XGBoost + SHAP", "Trust Score from 7 forensic signals", "88% acc · 0.95 AUC"),
        ("ManTraNet (CNN)", "Copy-move / splice localisation", "pixel-AUC 0.76–0.88"),
        ("Photo-region forensics", "Single-doc ID photo-swap", "96% recall · 0% FP"),
        ("FaceNet / MTCNN", "Cross-document face match", "96% match · 78% catch"),
        ("YOLOv8 ×5", "Stamp / Aadhaar / PAN / sig / utility", "stamps mAP 0.995"),
        ("spaCy trf + Verhoeff", "Names + PAN/Aadhaar/GSTIN/Udyam", "deterministic IDs"),
        ("Llama 3.1 8B (local)", "Plain-language evidence report", "on-device")]
cw, ch, cg = 5.95, 0.82, 0.12
for i, (m, role, met) in enumerate(rows):
    col, r = divmod(i, 4)
    x = 0.7 + col * (cw + 0.2); y = 2.45 + r * (ch + cg)
    card(s, x, y, cw, ch)
    text(s, x + 0.22, y + 0.1, 3.0, ch, m, size=13.5, color=TEXT, bold=True)
    text(s, x + 0.22, y + 0.43, 3.4, ch, role, size=10.5, color=MUTED)
    text(s, x + cw - 1.95, y + 0.18, 1.8, ch, met, size=11.5, color=NEON, bold=True, align=PP_ALIGN.RIGHT)

# ===================== 9 · FORENSICS =====================
s = slide(); header(s, "Forensic detection suite", "Ten ways a forged document gives itself away")
items = [("PDF metadata", "Producer/tool, edit traces, suspicious software."),
         ("Font & spacing", "Word-edit fingerprint: duplicated font subsets."),
         ("Error-Level Analysis", "Localised re-compression — pasted regions."),
         ("ManTraNet", "Copy-move & splicing heatmap (deep CNN)."),
         ("Stamp authentication", "YOLO + SIFT — the same stamp pasted twice."),
         ("Signature region", "ELA on the signature area."),
         ("Bank-statement balance", "Running-balance break = edited transaction."),
         ("Photo-region tamper", "Swapped ID photo (ManTraNet ∩ photo box)."),
         ("Cross-source mismatch", "Name vs DigiLocker/GSTN/Udyam record."),
         ("Fabricated Aadhaar", "Fails the UIDAI Verhoeff checksum.")]
cw, ch = 5.95, 0.78
for i, (t, d) in enumerate(items):
    col, r = divmod(i, 5)
    x = 0.7 + col * (cw + 0.2); y = 2.4 + r * (ch + 0.1)
    card(s, x, y, cw, ch)
    text(s, x + 0.22, y + 0.12, cw - 0.4, 0.4, t, size=13.5, color=NEON, bold=True)
    text(s, x + 0.22, y + 0.45, cw - 0.4, 0.4, d, size=11, color=MUTED)

# ===================== 10 · IDENTITY =====================
s = slide(); header(s, "Identity & cross-document", "Catching identity fraud a reviewer would miss")
bullets(s, 0.75, 2.45, 7.3, [
    ("Authoritative cross-source verification", "DigiLocker (PAN/Aadhaar), GSTN, Account Aggregator, Income-Tax e-Filing, Udyam — name must match the record."),
    ("Case-level consistency", "Same applicant across every document: conflicting hard-ID → RED; name/face mismatch → review."),
    ("Cross-document face match", "FaceNet across a case's IDs catches a mixed-identity photo."),
    ("Udyam QR authentication", "Decodes the certificate QR and cross-checks the URN."),
    ("Indian-name aware matching", "Fuzzy + phonetic (metaphone) so spelling variants still match."),
], gap=0.86)
card(s, 8.6, 2.45, 4.0, 3.9)
text(s, 8.85, 2.7, 3.5, 0.4, "Example caught", size=13, color=NEON, bold=True)
text(s, 8.85, 3.2, 3.5, 3.0, [
    [("Real PAN number,", TEXT, True, 14)], [("wrong name", RED, True, 14)], [("", MUTED, False, 6)],
    [("DigiLocker says “Rahul Verma”", MUTED, False, 12.5)], [("Document says “Suresh K”", MUTED, False, 12.5)],
    [("", MUTED, False, 6)], [("→ RED · fraud escalation", NEON, True, 14)],
], spacing=1.1)

# ===================== 11 · TIERS =====================
s = slide(); header(s, "Risk routing", "Three tiers — bank-safe, never auto-reject")
tiers = [("GREEN", "Fast-track", "Trust ≥ 85 and no critical indicator. Clean, system-generated, verified.", NEON),
         ("YELLOW", "Underwriter review", "Uncertain or an unreliable-on-real-docs flag. A human decides — never auto-rejected.", AMBER),
         ("RED", "Fraud escalation", "Trust < 50 or a critical forgery indicator (mismatch, tamper, fabricated ID).", RED)]
for i, (t, sub, d, c) in enumerate(tiers):
    x = 0.75 + i * 4.05
    card(s, x, 2.6, 3.85, 3.1)
    rect(s, x, 2.6, 3.85, 0.13, fill=c, radius=0.0, shape=MSO_SHAPE.RECTANGLE)
    text(s, x + 0.3, 2.95, 3.25, 0.6, t, size=26, color=c, bold=True)
    text(s, x + 0.3, 3.65, 3.25, 0.4, sub, size=15, color=TEXT, bold=True)
    text(s, x + 0.3, 4.15, 3.25, 1.4, d, size=12.5, color=MUTED, spacing=1.12)
text(s, 0.75, 6.05, 12, 0.6,
     [("Calibration favours false positives over false negatives ", TEXT, True, 14),
      ("— a missed forgery costs far more than a second look.", MUTED, False, 14)])

# ===================== 12 · COVERAGE =====================
s = slide(); header(s, "Document coverage", "26 document types across every loan-file category")
groups = [("Identity", ["Aadhaar", "PAN", "Passport"]),
          ("Financial", ["Bank statement", "ITR (full / V)", "Form-16", "Balance sheet", "P&L", "Cash-flow", "Audited financials", "Salary slip"]),
          ("Tax", ["GSTR-1", "GSTR-3B"]),
          ("Corporate / Legal", ["MOA / AOA", "Partnership deed", "Board resolution", "Power of attorney", "Sanction letter", "Guarantee letter", "Indemnity bond", "NOC", "Loan agreement", "Rental / lease"]),
          ("KYC / MSME", ["Udyam certificate", "Utility bill"])]
for i, (g, items) in enumerate(groups):
    x = 0.7 + i * 2.55
    card(s, x, 2.45, 2.42, 4.1)
    text(s, x + 0.2, 2.65, 2.02, 0.4, g, size=14, color=NEON, bold=True)
    text(s, x + 0.2, 3.17, 2.02, 3.2, [[("• " + it, TEXT, False, 12)] for it in items], spacing=1.15)
text(s, 0.72, 6.75, 12, 0.4, [("Bold additions this version: ", MUTED, False, 12),
     ("rental/lease · Udyam (MSME) + QR · utility bills (water/gas/electricity)", NEON, True, 12)])

# ===================== 13 · PERFORMANCE =====================
s = slide(); header(s, "Performance & validation", "Measured on a reproducible, labelled benchmark")
for i, (v, l, c) in enumerate([("98.6%", "Fraud hard-catch (332-doc benchmark)", NEON),
                               ("97.9%", "Doc-type accuracy (real held-out)", TEXT),
                               ("0.95", "XGBoost ROC-AUC (real+synthetic)", NEON),
                               ("91.6%", "Tamper recall (fraud)", TEXT)]):
    metric(s, 0.75 + i * 3.0, 2.5, 2.85, 1.5, v, l, vcolor=c)
card(s, 0.75, 4.2, 11.85, 2.0)
text(s, 1.0, 4.4, 11.3, 0.4, "data/benchmark.py — one harness, every change gated against it", size=13, color=NEON, bold=True)
text(s, 1.0, 4.9, 11.3, 1.3, [
    [("▸  ", NEON, True, 13), ("Two tracks (document pipeline vs image forensics) + 4 fraud vectors: metadata-tamper, splice, photo-swap, copy-paste.", TEXT, False, 13)],
    [("▸  ", NEON, True, 13), ("Reports precision / recall / false-positive rate, per-vector recall and a per-type confusion matrix.", TEXT, False, 13)],
    [("▸  ", NEON, True, 13), ("Honest: synthetic-only numbers are optimistic — these real+benchmark figures are the representative ones.", MUTED, False, 13)],
], spacing=1.15)

# ===================== 14 · FINAL-VERSION HARDENING =====================
s = slide(); header(s, "Final-version hardening", "What the benchmark drove us to fix")
bullets(s, 0.75, 2.45, 11.6, [
    ("Eliminated a false-rejection on genuine documents", "A genuine annual report was flagged “fabricated Aadhaar” (a tabular number); now context-gated — fixed, real fraud still caught."),
    ("Re-calibrated the photo-swap detector on real data", "In-box threshold 0.30→0.55 — genuine false-positives 7% → 0% with 94% swap recall."),
    ("Proved ManTraNet is correct on real PDFs", "The scary 26% synthetic false-positive was a render artifact; real born-digital PDFs trip it 0/24 — no calibration change needed."),
    ("Measured, then declined, a risky retrain", "Folding localised tampers into the whole-image scorer dropped fraud recall — reverted; localised detectors own that job."),
], gap=1.0)

# ===================== 15 · SECURITY / EXPLAINABILITY =====================
s = slide(); header(s, "Explainable · Secure · Compliant", "Built for a regulated bank")
cards = [("Explainable", ["SHAP per-feature attribution", "ELA / ManTraNet heatmaps", "Plain-language LLM report", "Dual-classifier agreement"]),
         ("Secure & on-prem", ["100% on-premise — no data egress", "OAuth2 + JWT, role-based access", "DPDP-Act / RBI friendly", "Cloud LLM optional drop-in"]),
         ("Compliant & auditable", ["Append-only audit log", "Every decision timestamped", "Immutable — never edited/deleted", "Full evidence frozen per decision"])]
for i, (t, items) in enumerate(cards):
    x = 0.75 + i * 4.05
    card(s, x, 2.55, 3.85, 3.5)
    text(s, x + 0.28, 2.8, 3.35, 0.5, t, size=16, color=NEON, bold=True)
    text(s, x + 0.28, 3.4, 3.35, 2.5, [[("✓  ", NEON, True, 13), (it, TEXT, False, 13)] for it in items], spacing=1.3)

# ===================== 16 · FLYWHEEL / CLOSE =====================
s = slide(); header(s, "Gets better in production", "The active-learning flywheel")
flow = ["Bank analyses\ndocuments", "Underwriter confirms\ngenuine / fraud", "Labels → real\ntraining manifest", "Models retrain\non real data", "Sharper detection\nnext time"]
bw, by = 2.18, 2.7
gap = (12.0 - bw) / (len(flow) - 1)
for i, st in enumerate(flow):
    x = 0.7 + i * gap
    card(s, x, by, bw, 1.45, fill=CARD2 if i % 2 else CARD)
    text(s, x + 0.1, by + 0.18, bw - 0.2, 1.1, [[(l, NEON if i in (1, 4) else TEXT, i in (1, 4), 12.5)] for l in st.split("\n")],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1.0)
    if i < len(flow) - 1:
        arrow(s, x + bw - 0.02, by + 0.5, gap - bw + 0.05)
text(s, 0.7, 4.7, 12, 0.5, "TrustLens", size=30, color=TEXT, bold=True)
text(s, 0.7, 5.35, 12, 0.5, [("Seconds, not days. Explainable, on-premise, and bank-safe.", NEON, True, 17)])
text(s, 0.7, 5.95, 12, 0.5, "98.6% fraud caught · 26 document types · 0 customer-data egress · never auto-rejects", size=13, color=MUTED)

import os
out = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "TrustLens_Final.pptx"))
prs.save(out)
print("Saved:", out, "·", len(prs.slides._sldIdLst), "slides")
